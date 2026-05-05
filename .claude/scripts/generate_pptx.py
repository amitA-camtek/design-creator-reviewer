"""
Static PowerPoint generator script for design packages.
Called by the powerpoint-generator skill with a slides-manifest.json path.
Usage: python generate_pptx.py <path-to-slides-manifest.json>
"""
import json
import sys
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("Installing python-pptx...")
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx", "-q"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN


def add_title_slide(prs, service_name, date_str):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = service_name
    tf = title.text_frame.paragraphs[0]
    tf.runs[0].font.size = Pt(36)
    tf.runs[0].font.bold = True
    tf.runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    subtitle.text = f"Design Package  |  {date_str}"
    subtitle.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)


def add_content_slide(prs, title_text, content_text):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].runs[0].font.size = Pt(28)
    title.text_frame.paragraphs[0].runs[0].font.bold = True

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, line in enumerate(content_text.splitlines()):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
    return slide


def main():
    if len(sys.argv) < 2:
        print("Usage: generate_pptx.py <slides-manifest.json>")
        sys.exit(1)

    manifest_path = sys.argv[1]
    with open(manifest_path, "r", encoding="utf-8") as f:
        manifest = json.load(f)

    service_name = manifest.get("service_name", "Service")
    output_path = manifest.get("output_path", "design.pptx")
    slides = manifest.get("slides", [])
    date_str = datetime.now().strftime("%Y-%m-%d")

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs, service_name, date_str)

    for slide_def in slides:
        add_content_slide(prs, slide_def.get("title", ""), slide_def.get("content", ""))

    prs.save(output_path)
    print(f"Saved: {output_path}")


if __name__ == "__main__":
    main()
